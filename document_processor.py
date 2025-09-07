import pandas as pd
import json
import xml.etree.ElementTree as ET
from docx import Document
from pptx import Presentation
import pdfplumber
from typing import List
from models import DataChunk
from text_normalizer import TextNormalizer
import io

class DocumentProcessor:
    def __init__(self, chunk_size: int = 5):
        self.chunk_size = chunk_size
        self.text_normalizer = TextNormalizer()
    
    def process_file(self, file_content: bytes, user_id: str, file_name: str) -> List[DataChunk]:
        file_ext = file_name.lower().split('.')[-1]
        
        if file_ext == 'csv':
            return self._process_csv(file_content, user_id, file_name)
        elif file_ext == 'json':
            return self._process_json(file_content, user_id, file_name)
        elif file_ext == 'xml':
            return self._process_xml(file_content, user_id, file_name)
        elif file_ext == 'txt':
            return self._process_text(file_content, user_id, file_name)
        elif file_ext in ['doc', 'docx']:
            return self._process_docx(file_content, user_id, file_name)
        elif file_ext in ['ppt', 'pptx']:
            return self._process_pptx(file_content, user_id, file_name)
        elif file_ext == 'pdf':
            return self._process_pdf(file_content, user_id, file_name)
        elif file_ext in ['xls', 'xlsx']:
            return self._process_excel(file_content, user_id, file_name)
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")
    
    def _process_csv(self, content: bytes, user_id: str, file_name: str) -> List[DataChunk]:
        df = pd.read_csv(io.BytesIO(content))
        chunks = []
        
        for i in range(0, len(df), self.chunk_size):
            chunk_df = df.iloc[i:i+self.chunk_size]
            text = self._dataframe_to_text(chunk_df)
            
            chunks.append(DataChunk(
                content=text,
                metadata={"file_name": file_name, "chunk_index": i // self.chunk_size, "type": "csv"},
                user_id=user_id
            ))
        return chunks
    
    def _process_json(self, content: bytes, user_id: str, file_name: str) -> List[DataChunk]:
        data = json.loads(content.decode('utf-8'))
        text = json.dumps(data, indent=2)
        
        return [DataChunk(
            content=text,
            metadata={"file_name": file_name, "type": "json"},
            user_id=user_id
        )]
    
    def _process_xml(self, content: bytes, user_id: str, file_name: str) -> List[DataChunk]:
        root = ET.fromstring(content.decode('utf-8'))
        text = ET.tostring(root, encoding='unicode')
        
        return [DataChunk(
            content=text,
            metadata={"file_name": file_name, "type": "xml"},
            user_id=user_id
        )]
    
    def _process_text(self, content: bytes, user_id: str, file_name: str) -> List[DataChunk]:
        text = content.decode('utf-8')
        
        return [DataChunk(
            content=text,
            metadata={"file_name": file_name, "type": "text"},
            user_id=user_id
        )]
    
    def _process_docx(self, content: bytes, user_id: str, file_name: str) -> List[DataChunk]:
        doc = Document(io.BytesIO(content))
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        
        return [DataChunk(
            content=text,
            metadata={"file_name": file_name, "type": "docx"},
            user_id=user_id
        )]
    
    def _process_pptx(self, content: bytes, user_id: str, file_name: str) -> List[DataChunk]:
        prs = Presentation(io.BytesIO(content))
        text_parts = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        
        text = "\n".join(text_parts)
        
        return [DataChunk(
            content=text,
            metadata={"file_name": file_name, "type": "pptx"},
            user_id=user_id
        )]
    
    def _process_pdf(self, content: bytes, user_id: str, file_name: str) -> List[DataChunk]:
        text_parts = []
        
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for page in pdf.pages:
                text_parts.append(page.extract_text() or "")
        
        text = "\n".join(text_parts)
        
        return [DataChunk(
            content=text,
            metadata={"file_name": file_name, "type": "pdf"},
            user_id=user_id
        )]
    
    def _process_excel(self, content: bytes, user_id: str, file_name: str) -> List[DataChunk]:
        df = pd.read_excel(io.BytesIO(content))
        chunks = []
        
        for i in range(0, len(df), self.chunk_size):
            chunk_df = df.iloc[i:i+self.chunk_size]
            text = self._dataframe_to_text(chunk_df)
            
            chunks.append(DataChunk(
                content=text,
                metadata={"file_name": file_name, "chunk_index": i // self.chunk_size, "type": "excel"},
                user_id=user_id
            ))
        return chunks
    
    def _dataframe_to_text(self, df: pd.DataFrame) -> str:
        text_parts = []
        for _, row in df.iterrows():
            normalized_items = []
            for col, val in row.items():
                if pd.notna(val):
                    # Normalize both column name and value
                    normalized_col = self.text_normalizer.get_canonical_form(str(col))
                    normalized_val = self.text_normalizer.get_canonical_form(str(val))
                    normalized_items.append(f"{normalized_col}: {normalized_val}")
            
            row_text = ", ".join(normalized_items)
            text_parts.append(row_text)
        return "\n".join(text_parts)